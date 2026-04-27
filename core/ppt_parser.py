"""
PPT 解析模块 - 将 PPTX 转换为图片

支持多平台自动选择最佳方案:
- Windows: pywin32 (PowerPoint COM) - 效果最好
- Linux/macOS: LibreOffice
- 保底: python-pptx 提取文本 + Pillow 生成占位图
"""

import os
import sys
import subprocess
from typing import List, Optional
from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PPTParser:
    def __init__(self, project_dir: str):
        self.project_dir = project_dir
        self.slides_dir = os.path.join(project_dir, "slides")

    def ppt_to_images(self, pptx_path: str) -> List[str]:
        """
        将 PPTX 转换为图片序列，自动选择最优方案
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPT 文件不存在: {pptx_path}")

        # 方案 1: Windows 下用 pywin32 (推荐)
        if sys.platform == "win32":
            images = self._try_win32(pptx_path)
            if images:
                return images

        # 方案 2: LibreOffice (跨平台)
        images = self._try_libreoffice(pptx_path)
        if images:
            return images

        # 方案 3: 保底，提取文本生成占位图
        print("[WARN] 无法渲染 PPT 为图片，使用文本提取模式")
        return self._extract_text_fallback(pptx_path)

    def _try_win32(self, pptx_path: str) -> Optional[List[str]]:
        """
        Windows 下使用 pywin32 调用 PowerPoint
        效果最好，可完美还原所有动画、排版、字体
        """
        try:
            import win32com.client
            import pythoncom
        except ImportError:
            print("[INFO] pywin32 未安装，跳过 Windows COM 方案")
            print("   如需使用，请: pip install pywin32")
            return None

        try:
            print("[WORK] 使用 Windows PowerPoint COM 渲染...")

            # 初始化 COM
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True  # 必须设为可见，否则部分版本会出错

            # 打开 PPT
            presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
            total_slides = presentation.Slides.Count

            # 导出每一页
            images = []
            for i in range(1, total_slides + 1):
                slide = presentation.Slides(i)
                filename = f"slide_{i:02d}.png"
                filepath = os.path.abspath(os.path.join(self.slides_dir, filename))
                slide.Export(filepath, "PNG", 1920, 1080)  # 1080P
                images.append(filename)
                print(f"   第 {i}/{total_slides} 页: {filename}")

            # 关闭 PPT
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()

            print(f"[OK] Windows COM 渲染完成: {len(images)} 页")
            return images

        except Exception as e:
            print(f"[WARN] Windows COM 方案出错: {e}")
            return None

    def _try_libreoffice(self, pptx_path: str) -> Optional[List[str]]:
        """
        使用 LibreOffice 转换 (跨平台方案)
        """
        try:
            # 检查 libreoffice 是否可用
            libreoffice_cmd = None
            for cmd in ["libreoffice", "soffice"]:
                try:
                    result = subprocess.run(
                        [cmd, "--version"], capture_output=True, text=True, timeout=5
                    )
                    if result.returncode == 0:
                        libreoffice_cmd = cmd
                        break
                except Exception:
                    continue

            if not libreoffice_cmd:
                return None

            print(f"[WORK] 使用 {libreoffice_cmd} 转换 PPT...")

            # 临时输出目录
            temp_dir = os.path.join(self.project_dir, "temp_ppt")
            os.makedirs(temp_dir, exist_ok=True)

            # 使用 LibreOffice 转换
            cmd = [
                libreoffice_cmd,
                "--headless",
                "--convert-to",
                "png",
                "--outdir",
                temp_dir,
                os.path.abspath(pptx_path),
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

            if result.returncode != 0:
                print(f"LibreOffice 转换失败: {result.stderr}")
                return None

            # 移动文件到 slides 目录
            images = []
            files = sorted(os.listdir(temp_dir))
            for i, filename in enumerate(files):
                if filename.lower().endswith(".png"):
                    new_name = f"slide_{i + 1:02d}.png"
                    src = os.path.join(temp_dir, filename)
                    dst = os.path.join(self.slides_dir, new_name)
                    os.rename(src, dst)
                    images.append(new_name)

            # 清理临时目录
            try:
                import shutil

                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception:
                pass

            if images:
                print(f"[OK] LibreOffice 转换完成: {len(images)} 页")
                return images
            return None

        except Exception as e:
            print(f"[WARN] LibreOffice 方案出错: {e}")
            return None

    def _extract_text_fallback(self, pptx_path: str) -> List[str]:
        """
        保底方案：提取文本，生成简单的占位图（改进版，避免文字重叠）
        """
        if not HAS_PPTX:
            print("[ERROR] python-pptx 也未安装，无法处理")
            return []

        prs = Presentation(pptx_path)
        images = []

        for i, slide in enumerate(prs.slides):
            # 提取文本
            text_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_lines.append(shape.text.strip())

            # 创建简单的占位图
            img = Image.new("RGB", (1920, 1080), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # 尝试加载字体（不同大小用于不同层级）
            font_title = None
            font_content = None
            font_small = None
            try_fonts = [
                "msyh.ttc",  # Windows 微软雅黑
                "simhei.ttf",  # Windows 黑体
                "PingFang.ttc",  # macOS 苹方
                "NotoSansCJK.ttc",  # Linux 思源黑体
                "DejaVuSans.ttf",  # Linux 通用
            ]
            for font_name in try_fonts:
                try:
                    font_title = ImageFont.truetype(font_name, 48)
                    font_content = ImageFont.truetype(font_name, 32)
                    font_small = ImageFont.truetype(font_name, 24)
                    break
                except Exception:
                    continue

            if not font_title:
                font_title = ImageFont.load_default()
                font_content = ImageFont.load_default()
                font_small = ImageFont.load_default()

            # 绘制标题
            page_num_text = f"第 {i + 1} 页"
            draw.text((100, 80), page_num_text, fill=(51, 51, 51), font=font_title)

            # 绘制分割线
            draw.line([(100, 160), (1820, 160)], fill=(200, 200, 200), width=2)

            # 绘制文本内容（智能换行，避免重叠）
            y = 200
            max_width = 1720  # 最大宽度，留出边距
            line_spacing = 45  # 行间距

            for line in text_lines:
                if not line:
                    continue

                # 如果是标题行（短文本），用粗体/更大字
                if len(line) <= 20 and y < 300:
                    current_font = font_title
                    current_spacing = 60
                    color = (30, 30, 30)
                else:
                    current_font = font_content
                    current_spacing = line_spacing
                    color = (60, 60, 60)

                # 智能换行：根据字体和宽度分割文本
                words = line.split()
                current_line = ""

                for word in words:
                    test_line = current_line + " " + word if current_line else word
                    bbox = draw.textbbox((0, 0), test_line, font=current_font)
                    test_width = bbox[2] - bbox[0]

                    if test_width <= max_width:
                        current_line = test_line
                    else:
                        if current_line:  # 画当前行
                            draw.text(
                                (100, y), current_line, fill=color, font=current_font
                            )
                            y += current_spacing
                        current_line = word

                if current_line:
                    draw.text((100, y), current_line, fill=color, font=current_font)
                    y += current_spacing

                # 如果超出范围，停止绘制
                if y > 1000:
                    break

            # 保存
            filename = f"slide_{i + 1:02d}.png"
            filepath = os.path.join(self.slides_dir, filename)
            img.save(filepath)
            images.append(filename)

        print(f"[OK] 文本提取完成: {len(images)} 页")
        print(f"[INFO] 建议安装 pywin32/LibreOffice 获得更好渲染效果")
        return images

    def get_slide_text(self, pptx_path: str) -> List[str]:
        """仅提取 PPT 每页的文本（用于生成讲解稿提示）"""
        if not HAS_PPTX:
            return []

        prs = Presentation(pptx_path)
        texts = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            texts.append("\n".join(slide_text))
        return texts
