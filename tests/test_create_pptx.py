#!/usr/bin/env python3
"""
创建一个测试用的 PPTX
"""
from pptx import Presentation
import os

prs = Presentation()

# 第 1 页
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "欢迎使用 CourseVideoGen"

# 第 2 页
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "什么是 CourseVideoGen？"
subtitle = slide.shapes.placeholders[1]
subtitle.text = "这是一个自动生成 PPT 讲解视频的工具。"

# 第 3 页
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "核心功能"
subtitle = slide.shapes.placeholders[1]
subtitle.text = "1. PPT 解析\n2. TTS 音频生成\n3. 视频合成"

output_path = "test.pptx"
prs.save(output_path)
print(f"✅ 测试 PPT 已生成: {os.path.abspath(output_path)}")

